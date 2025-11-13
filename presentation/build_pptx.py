"""
Build PowerPoint Presentation for Spotify Track Popularity Prediction

This script generates a professional PowerPoint presentation using captured
screenshots and model metadata.
"""

import json
from pathlib import Path
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configuration
SCREENS_DIR = Path(__file__).parent / "screens"
OUTPUT_DIR = Path(__file__).parent / "output"
TEMPLATE_DIR = Path(__file__).parent / "templates"
ASSETS_DIR = Path(__file__).parent.parent / "assets"
OUTPUT_FILE = OUTPUT_DIR / "Spotify_Popularity_Prediction_Presentation.pptx"
LOGO_PATH = ASSETS_DIR / "teamlogo.png"

# Team Color Theme (from teamlogo.png)
TEAM_NAVY = RGBColor(25, 42, 66)      # Dark navy blue background
TEAM_ORANGE = RGBColor(255, 100, 50)  # Orange accent
TEAM_WHITE = RGBColor(255, 255, 255)  # White text
TEAM_LIGHT_GRAY = RGBColor(240, 240, 240)  # Light gray for backgrounds

# Ensure output directory exists
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Slide content configuration
SLIDE_TITLES = {
    "01_xgboost_learning_curve.png": "Learning Curves - Model Training Progress",
    "02_actual_vs_predicted.png": "Model Performance - Actual vs Predicted",
    "03_prediction_density.png": "Prediction Density Distribution",
    "04_residuals_plot.png": "Residuals Analysis",
    "05_qq_plot_residuals.png": "QQ Plot - Residuals Normality Check",
    "06_correlation_heatmap.png": "Feature Correlation Heatmap",
    "07_feature_importance.png": "XGBoost Feature Importance",
    "08_xgboost_shap_summary_bar.png": "SHAP Values - Feature Impact (Bar)",
    "09_xgboost_shap_beeswarm.png": "SHAP Values - Feature Impact (Beeswarm)",
}

SLIDE_DESCRIPTIONS = {
    "01_xgboost_learning_curve.png": "Training and validation RMSE across iterations. No significant overfitting observed.",
    "02_actual_vs_predicted.png": "Strong correlation between actual and predicted popularity scores.",
    "03_prediction_density.png": "Density plot showing concentration of predictions around actual values.",
    "04_residuals_plot.png": "Residuals evenly distributed around zero, indicating good model fit.",
    "05_qq_plot_residuals.png": "Residuals follow normal distribution, validating model assumptions.",
    "06_correlation_heatmap.png": "Feature correlations reveal key relationships driving popularity.",
    "07_feature_importance.png": "XGBoost gain-based importance showing top predictive features.",
    "08_xgboost_shap_summary_bar.png": "SHAP values quantify each feature's average impact on predictions.",
    "09_xgboost_shap_beeswarm.png": "SHAP beeswarm plot shows feature impact direction and magnitude.",
}


def load_metadata():
    """Load metadata from capture_screens.py output"""
    metadata_file = SCREENS_DIR / "metadata.json"
    if metadata_file.exists():
        with open(metadata_file, 'r') as f:
            return json.load(f)
    return {}


def calculate_image_size(image_path, max_width, max_height):
    """
    Calculate image dimensions to fit within max bounds while preserving aspect ratio

    Args:
        image_path: Path to image file
        max_width: Maximum width in Inches
        max_height: Maximum height in Inches

    Returns:
        tuple: (width, height) in Inches
    """
    with Image.open(image_path) as img:
        img_width, img_height = img.size
        aspect_ratio = img_width / img_height

    # Calculate dimensions based on max constraints
    if aspect_ratio > max_width.inches / max_height.inches:
        # Image is wider - constrain by width
        width = max_width
        height = Inches(max_width.inches / aspect_ratio)
    else:
        # Image is taller - constrain by height
        height = max_height
        width = Inches(max_height.inches * aspect_ratio)

    return width, height


def create_title_slide(prs, metadata):
    """Create title slide with project metadata and team branding"""
    blank_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(blank_layout)

    # Add navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TEAM_NAVY

    # Add team logo in top-right corner
    if LOGO_PATH.exists():
        logo_width = Inches(1.5)
        logo_height = Inches(1.0)
        logo_left = Inches(8.3)
        logo_top = Inches(0.3)
        slide.shapes.add_picture(str(LOGO_PATH), logo_left, logo_top, width=logo_width)

    # Title
    title_left = Inches(0.8)
    title_top = Inches(2.5)
    title_width = Inches(8.4)
    title_height = Inches(1.2)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = metadata.get('title', 'Spotify Track Popularity Prediction')

    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEAM_WHITE

    # Subtitle with metrics
    subtitle_top = Inches(3.8)
    subtitle_height = Inches(2.0)

    subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True

    # Subtitle text
    subtitle_text = metadata.get('subtitle', 'Machine Learning Pipeline with XGBoost & SHAP')
    p1 = subtitle_frame.paragraphs[0]
    p1.text = subtitle_text
    p1.font.size = Pt(24)
    p1.font.color.rgb = TEAM_WHITE
    p1.space_after = Pt(18)

    # Add metrics if available
    if 'test_r2' in metadata:
        p2 = subtitle_frame.add_paragraph()
        p2.text = f"Test R²: {metadata['test_r2']:.4f}"
        if 'test_rmse' in metadata:
            p2.text += f"  |  RMSE: {metadata['test_rmse']:.2f}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = TEAM_ORANGE
        p2.font.bold = True
        p2.space_after = Pt(18)

    # Date
    p3 = subtitle_frame.add_paragraph()
    p3.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(200, 200, 200)


def create_agenda_slide(prs):
    """Create agenda/outline slide with team branding"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add team logo
    if LOGO_PATH.exists():
        logo_width = Inches(1.0)
        logo_left = Inches(8.8)
        logo_top = Inches(0.2)
        slide.shapes.add_picture(str(LOGO_PATH), logo_left, logo_top, width=logo_width)

    # Title with orange accent bar
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(8)
    title_height = Inches(0.7)

    # Orange accent bar
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.45),
        Inches(0.15), Inches(0.7)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAM_ORANGE
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(title_left + Inches(0.3), title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Presentation Agenda"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAM_NAVY

    # Add content text box
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    agenda_items = [
        "1. Learning Curves - Model Training Progress",
        "2. Model Performance Metrics",
        "3. Prediction Quality Analysis",
        "4. Residuals & Statistical Validation",
        "5. Feature Correlations",
        "6. Feature Importance Analysis",
        "7. SHAP Explainability",
        "8. Key Insights & Findings"
    ]

    for i, item in enumerate(agenda_items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = TEAM_NAVY
        p.space_after = Pt(14)


def create_content_slide(prs, image_path, title, description):
    """Create a content slide with properly sized image and team branding"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add team logo in top-right corner
    if LOGO_PATH.exists():
        logo_width = Inches(0.8)
        logo_left = Inches(9.0)
        logo_top = Inches(0.15)
        slide.shapes.add_picture(str(LOGO_PATH), logo_left, logo_top, width=logo_width)

    # Add title with orange accent bar
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(8.2)
    title_height = Inches(0.6)

    # Orange accent bar
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.28),
        Inches(0.1), Inches(0.6)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAM_ORANGE
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(title_left + Inches(0.2), title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title

    p = title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEAM_NAVY

    # Calculate proper image dimensions to fit in available space
    # Available space: from after title (1.1") to before description (6.3")
    # Available space height: 5.2 inches
    # Available width: 9 inches (with margins)
    max_img_width = Inches(9.0)
    max_img_height = Inches(4.8)

    img_width, img_height = calculate_image_size(image_path, max_img_width, max_img_height)

    # Center the image horizontally
    img_left = Inches(0.5 + (9.0 - img_width.inches) / 2)
    img_top = Inches(1.2)

    slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width, height=img_height)

    # Add description at bottom with light gray background
    desc_left = Inches(0.5)
    desc_top = Inches(6.5)
    desc_width = Inches(9)
    desc_height = Inches(0.9)

    # Light background box for description
    desc_bg = slide.shapes.add_shape(
        1,  # Rectangle
        desc_left, desc_top,
        desc_width, desc_height
    )
    desc_bg.fill.solid()
    desc_bg.fill.fore_color.rgb = TEAM_LIGHT_GRAY
    desc_bg.line.fill.background()

    desc_box = slide.shapes.add_textbox(desc_left + Inches(0.2), desc_top + Inches(0.1), desc_width - Inches(0.4), desc_height - Inches(0.2))
    desc_frame = desc_box.text_frame
    desc_frame.text = description
    desc_frame.word_wrap = True

    p = desc_frame.paragraphs[0]
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = RGBColor(50, 50, 50)


def create_summary_slide(prs, metadata):
    """Create final summary slide with team branding"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Add team logo
    if LOGO_PATH.exists():
        logo_width = Inches(1.0)
        logo_left = Inches(8.8)
        logo_top = Inches(0.2)
        slide.shapes.add_picture(str(LOGO_PATH), logo_left, logo_top, width=logo_width)

    # Title with orange accent bar
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(8)
    title_height = Inches(0.7)

    # Orange accent bar
    accent_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(0.45),
        Inches(0.15), Inches(0.7)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TEAM_ORANGE
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(title_left + Inches(0.3), title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = "Key Findings & Insights"

    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEAM_NAVY

    # Add content
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8.5)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    findings = [
        f"✅ Achieved R² of {metadata.get('test_r2', 0):.4f} on test set",
        f"✅ RMSE of {metadata.get('test_rmse', 0):.2f} demonstrates strong predictive accuracy",
        "✅ No significant overfitting observed in learning curves",
        "✅ Residuals follow normal distribution, validating model assumptions",
        "✅ SHAP analysis reveals key drivers of track popularity",
        f"✅ Model uses {len(metadata.get('features', []))} audio features for prediction",
        "✅ Full MLflow experiment tracking enables reproducibility",
    ]

    for i, finding in enumerate(findings):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = finding
        p.font.size = Pt(18)
        p.font.color.rgb = TEAM_NAVY
        p.space_after = Pt(16)


def build_presentation(metadata):
    """Build the complete presentation"""
    print("="*80)
    print("🎬 BUILDING POWERPOINT PRESENTATION")
    print("="*80)

    # Create presentation
    prs = Presentation()

    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\n📄 Adding slides...")

    # Slide 1: Title
    print("  1. Title slide")
    create_title_slide(prs, metadata)

    # Slide 2: Agenda
    print("  2. Agenda slide")
    create_agenda_slide(prs)

    # Slides 3-N: Content slides with images
    screenshots = sorted(SCREENS_DIR.glob("*.png"))
    slide_num = 3

    for screenshot in screenshots:
        if screenshot.name == "metadata.json":
            continue

        title = SLIDE_TITLES.get(screenshot.name, screenshot.stem.replace('_', ' ').title())
        description = SLIDE_DESCRIPTIONS.get(screenshot.name, "")

        print(f"  {slide_num}. {title}")
        create_content_slide(prs, screenshot, title, description)
        slide_num += 1

    # Final slide: Summary
    print(f"  {slide_num}. Summary slide")
    create_summary_slide(prs, metadata)

    # Save presentation
    prs.save(str(OUTPUT_FILE))

    print("\n" + "="*80)
    print("✅ PRESENTATION BUILT SUCCESSFULLY")
    print("="*80)
    print(f"\n📁 Output file: {OUTPUT_FILE}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📸 Screenshots used: {len(screenshots)}")

    return OUTPUT_FILE


def main():
    """Main execution"""

    # Load metadata
    metadata = load_metadata()

    # Check if screenshots exist
    screenshots = list(SCREENS_DIR.glob("*.png"))
    if not screenshots:
        print("❌ Error: No screenshots found in screens/ directory")
        print("   Please run 'python presentation/capture_screens.py' first")
        return

    # Build presentation
    output_file = build_presentation(metadata)

    print("\n🎯 Next steps:")
    print("  - Review presentation: open presentation/output/Spotify_Popularity_Prediction_Presentation.pptx")
    print("  - Customize slides as needed")
    print("  - Share with stakeholders")
    print("\n" + "="*80)


if __name__ == "__main__":
    main()
