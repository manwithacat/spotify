"""
Update PowerPoint presentation with new tuned model metrics
"""
from pptx import Presentation
from pptx.util import Pt
import re

# Load the presentation
prs = Presentation('presentation/Final_Spotify_Popularity_Prediction_Presentation.pptx')

# New metrics from hyperparameter tuning
NEW_METRICS = {
    'r2': 0.4772,
    'adj_r2': 0.4769,
    'rmse': 16.06,
    'mae': 11.95,
    'n_estimators': 450,
    'max_depth': 10,
    'learning_rate': 0.0689
}

print("🔄 Updating presentation with new tuned model metrics...")
print(f"   R² = {NEW_METRICS['r2']:.4f}")
print(f"   RMSE = {NEW_METRICS['rmse']:.2f}")
print(f"   MAE = {NEW_METRICS['mae']:.2f}")
print(f"   Estimators = {NEW_METRICS['n_estimators']}\n")

updates_made = 0

for slide_num, slide in enumerate(prs.slides, 1):
    slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {slide_num}"

    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue

        original_text = shape.text
        modified = False

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text

                # Update R² values
                # Replace 0.8447, 0.39, 0.36, 0.25 with new R² 0.4772
                if re.search(r'R²\s*[=:]\s*0\.8447', text, re.IGNORECASE):
                    new_text = re.sub(r'(R²\s*[=:]\s*)0\.8447', rf'\g<1>{NEW_METRICS["r2"]:.4f}', text, flags=re.IGNORECASE)
                    run.text = new_text
                    modified = True
                    print(f"✓ Slide {slide_num} ({slide_title}): Updated R² = 0.8447 → {NEW_METRICS['r2']:.4f}")

                elif re.search(r'R²\s*[=:]\s*0\.39', text, re.IGNORECASE):
                    new_text = re.sub(r'(R²\s*[=:]\s*)0\.39', rf'\g<1>{NEW_METRICS["r2"]:.4f}', text, flags=re.IGNORECASE)
                    run.text = new_text
                    modified = True
                    print(f"✓ Slide {slide_num} ({slide_title}): Updated R² = 0.39 → {NEW_METRICS['r2']:.4f}")

                elif re.search(r'R²\s*[=:]\s*0\.36', text, re.IGNORECASE):
                    # This is Random Forest - leave as is for comparison
                    pass

                elif re.search(r'R²\s*[=:]\s*0\.25', text, re.IGNORECASE):
                    # This is Linear Regression - leave as is for comparison
                    pass

                # Update "best model" or "final model" R² references
                if 'xgboost' in text.lower() and re.search(r'0\.\d{2,4}', text):
                    # Replace any decimal that looks like an R² score after XGBoost mention
                    new_text = re.sub(r'\b0\.8447\b', f'{NEW_METRICS["r2"]:.4f}', text)
                    if new_text != text:
                        run.text = new_text
                        modified = True
                        print(f"✓ Slide {slide_num} ({slide_title}): Updated XGBoost R² in context")

                # Update RMSE values
                if re.search(r'RMSE\s*[=:]\s*5\.74', text, re.IGNORECASE):
                    new_text = re.sub(r'(RMSE\s*[=:]\s*)5\.74', rf'\g<1>{NEW_METRICS["rmse"]:.2f}', text, flags=re.IGNORECASE)
                    run.text = new_text
                    modified = True
                    print(f"✓ Slide {slide_num} ({slide_title}): Updated RMSE = 5.74 → {NEW_METRICS['rmse']:.2f}")

                # Update MAE values
                if re.search(r'MAE\s*[=:]\s*4\.54', text, re.IGNORECASE):
                    new_text = re.sub(r'(MAE\s*[=:]\s*)4\.54', rf'\g<1>{NEW_METRICS["mae"]:.2f}', text, flags=re.IGNORECASE)
                    run.text = new_text
                    modified = True
                    print(f"✓ Slide {slide_num} ({slide_title}): Updated MAE = 4.54 → {NEW_METRICS['mae']:.2f}")

                # Update n_estimators
                if re.search(r'200\s+estimators', text, re.IGNORECASE):
                    new_text = re.sub(r'200(\s+estimators)', rf'{NEW_METRICS["n_estimators"]}\g<1>', text, flags=re.IGNORECASE)
                    run.text = new_text
                    modified = True
                    print(f"✓ Slide {slide_num} ({slide_title}): Updated estimators = 200 → {NEW_METRICS['n_estimators']}")

                if modified:
                    updates_made += 1

# Save the updated presentation
output_path = 'presentation/Final_Spotify_Popularity_Prediction_Presentation_Updated.pptx'
prs.save(output_path)

print(f"\n✅ Presentation updated successfully!")
print(f"   Total updates: {updates_made}")
print(f"   Saved to: {output_path}")
print(f"\n📝 Summary of changes:")
print(f"   • R² = 0.4772 (was 0.8447 or 0.39)")
print(f"   • RMSE = 16.06 (was 5.74 or similar)")
print(f"   • MAE = 11.95 (was 4.54 or similar)")
print(f"   • Estimators = 450 (was 200)")
print(f"\n💡 Note: Model comparison R² values (Linear=0.25, RandomForest=0.36) kept for comparison")
